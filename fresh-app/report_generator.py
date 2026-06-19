#!/usr/bin/env python3
"""
Report generator — called by Deno via subprocess.
Reads JSON payload from stdin, writes binary (PDF or PPTX) to stdout.
"""
import sys
import json
import base64
import io
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas as rl_canvas
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def generate_pdf(payload: dict) -> bytes:
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=landscape(A4))
    w, h = landscape(A4)

    # Título
    c.setFont("Helvetica-Bold", 20)
    c.drawCentredString(w / 2, h - 60, payload.get("title", "Relatório"))
    c.setFont("Helvetica", 12)
    c.drawCentredString(w / 2, h - 85, payload.get("subtitle", ""))

    # Gráficos
    for chart in payload.get("charts", []):
        c.showPage()
        c.setFont("Helvetica-Bold", 14)
        c.drawCentredString(w / 2, h - 40, chart.get("title", ""))
        img_data = chart.get("image", "")
        if "," in img_data:
            img_data = img_data.split(",")[1]
        try:
            img_bytes = base64.b64decode(img_data)
            img_buf = io.BytesIO(img_bytes)
            img_w, img_h = w - 80, h - 100
            c.drawImage(img_buf, 40, 60, width=img_w, height=img_h, preserveAspectRatio=True)
        except Exception:
            pass

    # Tabela de dados
    data = payload.get("data", [])
    headers = payload.get("headers", {})
    if data and headers:
        c.showPage()
        c.setFont("Helvetica-Bold", 14)
        c.drawCentredString(w / 2, h - 40, "Dados")
        keys = list(headers.keys())
        col_w = (w - 60) / max(len(keys), 1)
        y = h - 70
        c.setFont("Helvetica-Bold", 7)
        for i, k in enumerate(keys):
            c.drawCentredString(40 + i * col_w + col_w / 2, y, str(headers[k])[:20])
        y -= 14
        for row in data:
            if y < 40:
                c.showPage()
                y = h - 40
            c.setFont("Helvetica", 6)
            for i, k in enumerate(keys):
                c.drawCentredString(40 + i * col_w + col_w / 2, y, str(row.get(k, ""))[:25])
            y -= 10

    c.save()
    return buf.getvalue()


def generate_pptx(payload: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    # Slide 1: título
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    tf = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.3), Inches(1.5))
    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = payload.get("title", "Relatório")
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tf2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.8))
    p2 = tf2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = payload.get("subtitle", "")
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

    # Slides de gráficos
    for chart in payload.get("charts", []):
        slide = prs.slides.add_slide(blank)
        tf = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.5))
        p = tf.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = chart.get("title", "")
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

        img_data = chart.get("image", "")
        if "," in img_data:
            img_data = img_data.split(",")[1]
        try:
            img_bytes = base64.b64decode(img_data)
            img_buf = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_buf, Inches(0.3), Inches(0.7), Inches(12.7), Inches(6.5))
        except Exception:
            pass

    # Slide de tabela
    data = payload.get("data", [])
    headers = payload.get("headers", {})
    if data and headers:
        slide = prs.slides.add_slide(blank)
        tf = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.4))
        run = tf.text_frame.paragraphs[0].add_run()
        run.text = "Dados"
        run.font.size = Pt(14)
        run.font.bold = True

        keys = list(headers.keys())
        rows_data = data[:30]
        table = slide.shapes.add_table(len(rows_data) + 1, len(keys), Inches(0.2), Inches(0.6), Inches(12.9), Inches(6.5)).table
        col_w = Inches(12.9 / max(len(keys), 1))
        for col in table.columns:
            col.width = col_w

        for i, k in enumerate(keys):
            cell = table.cell(0, i)
            cell.text = str(headers[k])
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(7)

        for ri, row in enumerate(rows_data):
            for ci, k in enumerate(keys):
                cell = table.cell(ri + 1, ci)
                cell.text = str(row.get(k, ""))
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(6)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    payload = json.loads(sys.stdin.buffer.read())
    fmt = payload.get("format", "pdf")
    if fmt == "pptx":
        result = generate_pptx(payload)
    else:
        result = generate_pdf(payload)
    sys.stdout.buffer.write(result)
