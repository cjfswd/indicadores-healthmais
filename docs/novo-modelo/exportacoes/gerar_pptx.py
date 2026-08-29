# -*- coding: utf-8 -*-
from exp_common import *
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

C = lambda h: RGBColor.from_string(h)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def slide():
    return prs.slides.add_slide(BLANK)


def caixa(s, x, y, w, h, texto, tam=14, cor=INK, bold=False, align=PP_ALIGN.LEFT, espaco=1.15):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = espaco
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(tam)
    r.font.bold = bold
    r.font.name = "Calibri"
    r.font.color.rgb = C(cor)
    return tb


def retangulo(s, x, y, w, h, fill, linha=None):
    sh = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C(fill)
    if linha:
        sh.line.color.rgb = C(linha)
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.03
    return sh


# ── Capa ──────────────────────────────────────────────────
s = slide()
retangulo(s, -0.1, -0.1, W + 0.2, 7.7, PINE)
caixa(s, 1.1, 2.35, 10, 0.4, "HEALTHMAIS · ATENDIMENTO DOMICILIAR", 13, "BFD8CF", True)
caixa(s, 1.1, 2.82, 11, 1.1, "Painel de Indicadores", 46, "FFFFFF", True)
caixa(s, 1.1, 3.95, 11, 0.5, "Competência " + COMP + " · modelo recategorizado em dez cards", 17, "D8E8E1")
retangulo(s, 1.12, 4.8, 2.2, 0.045, "8FC0B2")
caixa(s, 1.1, 5.15, 10.5, 0.4,
      "Gerado em 28/08/2026 · recorte: todos os períodos, todas as operadoras, AD e ID", 11, "A9CCC0")

# ── Sumário ───────────────────────────────────────────────
s = slide()
caixa(s, 0.85, 0.6, 10, 0.5, "Os dez cards no período", 28, INK, True)
caixa(s, 0.85, 1.15, 11.6, 0.35, "Valor principal de cada indicador e situação frente à meta.", 13, INK2)
x0, y0, lw, lh, gx, gy = 0.85, 1.8, 2.85, 1.32, 0.14, 0.16
for i, c in enumerate(CARDS):
    col, lin = i % 4, i // 4
    x = x0 + col * (lw + gx)
    y = y0 + lin * (lh + gy)
    fora = c['meta'] and not c['meta']['ok']
    retangulo(s, x, y, lw, lh, BRICK_L if fora else SURF2, BRICK if fora else LINE)
    caixa(s, x + 0.18, y + 0.14, lw - 0.3, 0.25, c['code'], 10, BRICK if fora else PINE, True)
    caixa(s, x + 0.18, y + 0.38, lw - 0.3, 0.5, c['nome'], 11.5, INK, True, espaco=0.95)
    st = c['stats'][0]
    caixa(s, x + 0.18, y + 0.84, lw - 0.3, 0.35, st['v'], 19, BRICK if fora else PINE, True)
    caixa(s, x + 0.18, y + 1.11, lw - 0.3, 0.2, st['k'].lower(), 8.5, INK3)
caixa(s, 0.85, 6.95, 11.6, 0.3,
      "Card 03 é espelho de 1.4 e não soma ao painel · card 07 é campo de valor, não contador", 9, INK3)

# ── Um slide por card ─────────────────────────────────────
for c in CARDS:
    s = slide()
    retangulo(s, -0.1, -0.1, W + 0.2, 0.18, PINE)
    caixa(s, 0.85, 0.45, 1.5, 0.3, "CARD " + c['code'], 11, INK3, True)
    caixa(s, 0.85, 0.75, 8.6, 0.5, c['nome'], 26, INK, True)
    caixa(s, 0.85, 1.3, 8.4, 0.62, c['nota'], 11.5, INK2, espaco=1.2)

    if c['meta']:
        ok = c['meta']['ok']
        mw = 3.35
        retangulo(s, W - mw - 0.85, 0.7, mw, 0.9, PINE_L if ok else BRICK_L, PINE if ok else BRICK)
        caixa(s, W - mw - 0.68, 0.84, mw - 0.34, 0.25,
              "META " + ("ATINGIDA" if ok else "NÃO ATINGIDA"), 9.5, PINE if ok else BRICK, True)
        caixa(s, W - mw - 0.68, 1.08, mw - 0.34, 0.48, c['meta']['txt'], 10.5, INK2, espaco=1.05)

    kw = 2.85
    for i, st in enumerate(c['stats'][:4]):
        x = 0.85 + i * (kw + 0.14)
        retangulo(s, x, 2.12, kw, 1.16, BRICK_L if st['warn'] else (PINE_L if i == 0 else SURF2),
                  BRICK if st['warn'] else LINE)
        cor = BRICK if st['warn'] else (PINE if i == 0 else INK)
        caixa(s, x + 0.18, 2.25, kw - 0.3, 0.22, st['k'].upper(), 8, INK3, True)
        caixa(s, x + 0.18, 2.48, kw - 0.3, 0.4, st['v'], 21, cor, True)
        caixa(s, x + 0.18, 2.9, kw - 0.3, 0.3, st['d'], 9, INK3, espaco=0.95)

    piv = c['pivot']
    meses = piv['cols'][1:-1]
    pai = next((r for r in piv['rows'] if r['pai']), None)
    vals = [num(v) or 0 for v in pai['cells'][1:-1]] if pai else []
    caixa(s, 0.85, 3.5, 5.6, 0.25, "EVOLUÇÃO MENSAL", 8.5, INK3, True)
    if meses and vals:
        cd = CategoryChartData()
        cd.categories = meses
        cd.add_series("Total", vals)
        gr = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(3.72),
                                Inches(5.9), Inches(2.9), cd).chart
        gr.has_legend = False
        gr.plots[0].gap_width = 70
        ser = gr.plots[0].series[0]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = C(PINE)
        gr.value_axis.has_major_gridlines = True
        gr.value_axis.major_gridlines.format.line.color.rgb = C(LINE)
        gr.category_axis.tick_labels.font.size = Pt(10)
        gr.value_axis.tick_labels.font.size = Pt(10)

    caixa(s, 6.95, 3.5, 5.4, 0.25, "POR SUBCATEGORIA", 8.5, INK3, True)
    subs = [r for r in piv['rows'] if not r['pai']]
    subs = sorted(subs, key=lambda r: -(num(r['cells'][-1]) or 0))[:7]
    maxv = max([num(r['cells'][-1]) or 0 for r in subs] + [1])
    y = 3.8
    for r in subs:
        v = num(r['cells'][-1]) or 0
        caixa(s, 6.95, y, 4.0, 0.25, r['cells'][0][:50], 10.5, INK2 if v else INK3)
        larg = (v / maxv) * 1.0
        if larg > 0.02:
            retangulo(s, 11.0, y + 0.07, larg, 0.12, PINE)
        caixa(s, 12.05, y - 0.01, 1.0, 0.25, r['cells'][-1], 10.5, INK, True, PP_ALIGN.RIGHT)
        y += 0.39

    if c['alertas']:
        n = len(c['alertas'])
        retangulo(s, 0.85, 6.76, 11.6, 0.5, BRICK_L, BRICK)
        plural = "s" if n > 1 else ""
        txt = "%d pendência%s em aberto · %s" % (n, plural, c['alertas'][0]['texto'][:94])
        caixa(s, 1.05, 6.89, 11.2, 0.3, txt, 10, BRICK)

# ── Pendências ───────────────────────────────────────────
s = slide()
caixa(s, 0.85, 0.6, 10, 0.5, "Pendências em aberto", 28, INK, True)
todos = [(c['code'], a) for c in CARDS for a in c['alertas']]
caixa(s, 0.85, 1.15, 11.6, 0.35,
      "%d itens exigem ação antes do fechamento da competência." % len(todos), 13, INK2)
y = 1.78
for code, a in todos[:12]:
    retangulo(s, 0.85, y, 11.6, 0.4, SURF2, LINE)
    caixa(s, 1.02, y + 0.1, 0.5, 0.22, code, 9.5, PINE, True)
    caixa(s, 1.55, y + 0.1, 2.3, 0.22, a['quem'][:28], 10, INK, True)
    caixa(s, 3.95, y + 0.1, 8.3, 0.22, a['texto'][:112], 9.5, INK2)
    y += 0.46
if len(todos) > 12:
    caixa(s, 0.85, y + 0.06, 11.6, 0.3,
          "e mais %d pendências na aba Alertas da planilha." % (len(todos) - 12), 10, INK3)

# ── Retrocompatibilidade ──────────────────────────────────
s = slide()
caixa(s, 0.85, 0.6, 10, 0.5, "O que vem do sistema atual", 28, INK, True)
caixa(s, 0.85, 1.15, 11.6, 0.35, "Nenhuma função existente se perde na virada de modelo.", 13, INK2)
y = 1.85
for p in D['paginas'][:9]:
    retangulo(s, 0.85, y, 11.6, 0.52, SURF2, LINE)
    caixa(s, 1.05, y + 0.15, 2.8, 0.25, p['nome'], 11.5, INK, True)
    mantido = next((v for k, v in p['itens'] if k == "Mantido"), p['itens'][0][1])
    caixa(s, 3.95, y + 0.16, 8.3, 0.25, mantido[:118], 10, INK2)
    y += 0.58

prs.save("Painel_Indicadores_08-2026.pptx")
print("pptx ok -", len(prs.slides._sldIdLst), "slides")
